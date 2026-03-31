"""
Serviço de exportação de apresentações para PPTX.

Converte o JSON de SlidePresentation em arquivo .pptx
usando as cores e fontes do SlideTemplate associado.
"""

import io
import logging

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Converte cor hex (#RRGGBB) para RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def generate_pptx(
    title: str,
    slides_content: list[dict],
    template_rules: dict | None = None,
) -> bytes:
    """
    Gera um arquivo .pptx a partir do JSON de slides.

    Args:
        title: Título da apresentação
        slides_content: Lista de slides [{title, content, notes, layout}]
        template_rules: Regras do template (cores, fontes, layout)

    Returns:
        Bytes do arquivo .pptx
    """
    rules = template_rules or {}
    colors = rules.get("colors", {})
    fonts = rules.get("fonts", {})
    layout_rules = rules.get("layout", {})

    primary = colors.get("primary", "#1B3A5C")
    secondary = colors.get("secondary", "#E8B931")
    bg_color = colors.get("background", "#FFFFFF")
    text_color = colors.get("text", "#333333")
    title_font = fonts.get("title", "Calibri")
    body_font = fonts.get("body", "Calibri")
    title_size = fonts.get("size_title", 32)
    body_size = fonts.get("size_body", 18)
    footer_text = layout_rules.get("footer_text", "")
    show_numbers = layout_rules.get("include_slide_numbers", True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in slides_content:
        slide_layout = slide_data.get("layout", "default")
        slide_title = slide_data.get("title", "")
        content_items = slide_data.get("content", [])
        notes = slide_data.get("notes", "")
        slide_num = slide_data.get("slide_number", 0)

        # Usa layout em branco para controle total
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Fundo
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(bg_color)

        if slide_layout == "title":
            # Slide de capa
            _add_title_slide(slide, slide_title, primary, secondary, title_font, title_size)
        else:
            # Barra superior com cor primária
            _add_header_bar(slide, primary)

            # Título do slide
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(0.8))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_title
            p.font.size = Pt(title_size)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.name = title_font

            # Conteúdo (bullets)
            content_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(1.5), Inches(11.5), Inches(5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True

            for i, item in enumerate(content_items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"  {item}"
                p.font.size = Pt(body_size)
                p.font.color.rgb = _hex_to_rgb(text_color)
                p.font.name = body_font
                p.space_after = Pt(8)
                p.level = 0

        # Rodapé
        if footer_text or show_numbers:
            footer_parts = []
            if footer_text:
                footer_parts.append(footer_text)
            if show_numbers and slide_num:
                footer_parts.append(str(slide_num))
            _add_footer(slide, " | ".join(footer_parts), text_color, body_font)

        # Notas do professor
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_title_slide(
    slide,
    title: str,
    primary: str,
    secondary: str,
    font_name: str,
    font_size: int,
):
    """Cria slide de capa com cor primária de fundo."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(primary)

    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(font_size + 8)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.name = font_name
    p.alignment = PP_ALIGN.CENTER

    # Linha decorativa
    accent_box = slide.shapes.add_textbox(Inches(4), Inches(4.6), Inches(5), Inches(0.1))
    accent_fill = accent_box.fill
    accent_fill.solid()
    accent_fill.fore_color.rgb = _hex_to_rgb(secondary)


def _add_header_bar(slide, primary: str):
    """Adiciona barra de cabeçalho colorida."""
    from pptx.util import Emu

    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Emu(0),
        Emu(0),
        Inches(13.333),
        Inches(1.2),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(primary)
    shape.line.fill.background()


def _add_footer(slide, text: str, color: str, font_name: str):
    """Adiciona rodapé."""
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12), Inches(0.4))
    tf = footer_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = _hex_to_rgb(color)
    p.font.name = font_name
    p.alignment = PP_ALIGN.RIGHT
